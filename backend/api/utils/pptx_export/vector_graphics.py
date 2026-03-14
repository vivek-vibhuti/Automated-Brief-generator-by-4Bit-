"""
Vector graphics generation for PowerPoint slides using SVG and matplotlib
"""
import io
import base64
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import numpy as np
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
import tempfile
import os

class VectorGraphicsGenerator:
    """Generates vector graphics and charts for PowerPoint"""
    
    @staticmethod
    def create_bar_chart(labels, values, title="Data Visualization"):
        """Create a bar chart as EMF vector image"""
        try:
            # Create figure with vector format
            fig, ax = plt.subplots(figsize=(6, 4), dpi=100)
            
            # Create bars with Hawkins Lab colors
            bars = ax.bar(labels, values, color='#00ff00', alpha=0.7, edgecolor='#00aa00')
            
            # Style the chart
            ax.set_title(title, color='white', fontsize=14, fontweight='bold')
            ax.set_facecolor('#111111')
            fig.patch.set_facecolor('#222222')
            
            # Style axes
            ax.tick_params(colors='#00ff00')
            ax.spines['bottom'].set_color('#00ff00')
            ax.spines['left'].set_color('#00ff00')
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            
            # Add value labels on bars
            for bar in bars:
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height,
                       f'{height}', ha='center', va='bottom', color='white')
            
            # Save as vector format (EMF for PowerPoint)
            with tempfile.NamedTemporaryFile(suffix='.emf', delete=False) as tmp:
                plt.savefig(tmp.name, format='emf', bbox_inches='tight', 
                          facecolor=fig.get_facecolor(), edgecolor='none')
                plt.close()
                return tmp.name
                
        except Exception as e:
            print(f"Bar chart error: {e}")
            return None
    
    @staticmethod
    def create_pie_chart(values, labels, title="Distribution"):
        """Create a pie chart as vector image"""
        try:
            fig, ax = plt.subplots(figsize=(6, 4), dpi=100)
            
            # Hawkins Lab color scheme
            colors = ['#00ff00', '#ff0000', '#0000ff', '#ffff00', '#ff00ff', '#00ffff']
            
            # Create pie chart
            wedges, texts, autotexts = ax.pie(
                values, labels=labels, autopct='%1.1f%%',
                colors=colors[:len(values)],
                textprops={'color': 'white', 'fontweight': 'bold'}
            )
            
            # Style
            ax.set_title(title, color='white', fontsize=14, fontweight='bold')
            fig.patch.set_facecolor('#222222')
            
            # Save as vector format
            with tempfile.NamedTemporaryFile(suffix='.emf', delete=False) as tmp:
                plt.savefig(tmp.name, format='emf', bbox_inches='tight',
                          facecolor=fig.get_facecolor(), edgecolor='none')
                plt.close()
                return tmp.name
                
        except Exception as e:
            print(f"Pie chart error: {e}")
            return None
    
    @staticmethod
    def create_line_chart(x_values, y_values, title="Trend Analysis"):
        """Create a line chart as vector image"""
        try:
            fig, ax = plt.subplots(figsize=(6, 4), dpi=100)
            
            # Create line chart
            ax.plot(x_values, y_values, color='#00ff00', linewidth=2, marker='o')
            ax.fill_between(x_values, y_values, alpha=0.2, color='#00ff00')
            
            # Style
            ax.set_title(title, color='white', fontsize=14, fontweight='bold')
            ax.set_facecolor('#111111')
            fig.patch.set_facecolor('#222222')
            
            ax.tick_params(colors='#00ff00')
            ax.spines['bottom'].set_color('#00ff00')
            ax.spines['left'].set_color('#00ff00')
            
            # Save as vector format
            with tempfile.NamedTemporaryFile(suffix='.emf', delete=False) as tmp:
                plt.savefig(tmp.name, format='emf', bbox_inches='tight',
                          facecolor=fig.get_facecolor(), edgecolor='none')
                plt.close()
                return tmp.name
                
        except Exception as e:
            print(f"Line chart error: {e}")
            return None
    
    @staticmethod
    def add_vector_logo(slide, position=(1, 1), size=(1, 1)):
        """Add Hawkins Lab vector logo to slide"""
        try:
            # Create simple vector logo using shapes
            left = Inches(position[0])
            top = Inches(position[1])
            width = Inches(size[0])
            height = Inches(size[1])
            
            # Add circle (lab logo)
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left, top, width, height
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = None  # Transparent
            circle.line.color.rgb = None  # Green outline
            
            # Add text
            textbox = slide.shapes.add_textbox(
                left + Inches(0.1), top + Inches(0.1),
                width - Inches(0.2), height - Inches(0.2)
            )
            tf = textbox.text_frame
            tf.text = "HL"
            tf.paragraphs[0].font.color.rgb = None  # Green
            tf.paragraphs[0].font.size = Inches(0.3)
            tf.paragraphs[0].font.bold = True
            
            return True
        except Exception as e:
            print(f"Logo error: {e}")
            return False
    
    @staticmethod
    def create_timeline_diagram(events, dates, title="Project Timeline"):
        """Create a timeline diagram as vector"""
        try:
            fig, ax = plt.subplots(figsize=(8, 3), dpi=100)
            
            # Create timeline
            y_pos = np.zeros(len(dates))
            ax.scatter(range(len(dates)), y_pos, c='#00ff00', s=100)
            
            # Add event labels
            for i, (date, event) in enumerate(zip(dates, events)):
                ax.annotate(f"{date}\n{event}", (i, 0), 
                          xytext=(0, 30 if i % 2 == 0 else -30),
                          textcoords='offset points', ha='center',
                          color='white', fontsize=8,
                          bbox=dict(boxstyle='round,pad=0.3',
                                  facecolor='#333333', edgecolor='#00ff00'))
            
            # Style
            ax.set_title(title, color='white', fontsize=14, fontweight='bold')
            ax.set_facecolor('#111111')
            fig.patch.set_facecolor('#222222')
            ax.set_ylim(-1, 1)
            ax.axis('off')
            
            # Save as vector
            with tempfile.NamedTemporaryFile(suffix='.emf', delete=False) as tmp:
                plt.savefig(tmp.name, format='emf', bbox_inches='tight',
                          facecolor=fig.get_facecolor(), edgecolor='none')
                plt.close()
                return tmp.name
                
        except Exception as e:
            print(f"Timeline error: {e}")
            return None