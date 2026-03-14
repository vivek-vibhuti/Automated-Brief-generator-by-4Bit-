"""
Enhanced PowerPoint exporter with animations and vector graphics
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile
import os
from .animations import AnimationManager
from .vector_graphics import VectorGraphicsGenerator
import logging

logger = logging.getLogger(__name__)

class EnhancedPPTXExporter:
    """Creates professional PowerPoint presentations with animations and graphics"""
    
    def __init__(self):
        self.prs = Presentation()
        self.animations = AnimationManager()
        self.graphics = VectorGraphicsGenerator()
        self.temp_files = []  # Track temp files for cleanup
    
    def _cleanup_temp_files(self):
        """Clean up temporary files"""
        for file in self.temp_files:
            try:
                if os.path.exists(file):
                    os.unlink(file)
            except:
                pass
    
    def _add_title_slide(self, slide_data):
        """Add a title slide with animations"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        if title:
            title.text = slide_data.get('title', 'Hawkins Briefing')
            # Style title
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)  # Red
            title.text_frame.paragraphs[0].font.size = Pt(44)
            title.text_frame.paragraphs[0].font.bold = True
        
        # Set subtitle
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = "Mr. Clarke's Automated Briefing Generator"
        
        # Add animations
        self.animations.add_fade_animation(title, slide, timing_sec=0.5)
        
        return slide
    
    def _add_content_slide(self, slide_data, index):
        """Add a content slide with bullet animations"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        if title:
            title.text = slide_data.get('title', f'Slide {index + 1}')
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
            title.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Add bullet points
        content = slide.placeholders[1]
        if content:
            tf = content.text_frame
            tf.clear()
            
            bullets = slide_data.get('bullets', [])
            for i, bullet in enumerate(bullets):
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0, 255, 0)  # Green
                
                # First bullet bold
                if i == 0:
                    p.font.bold = True
            
            # Add sequential bullet animations
            self.animations.add_bullet_animations(slide, slide_data)
        
        # Add vector chart if present
        if slide_data.get('chart'):
            self._add_vector_chart(slide, slide_data['chart'])
        
        # Add slide transition
        self.animations.add_slide_transition(slide, transition_type='fade', duration=1.0)
        
        return slide
    
    def _add_vector_chart(self, slide, chart_data):
        """Add a vector chart to slide"""
        chart_path = None
        try:
            # Generate appropriate chart based on type
            chart_type = chart_data.get('type', 'bar')
            
            if chart_type == 'bar':
                chart_path = self.graphics.create_bar_chart(
                    chart_data.get('labels', ['A', 'B', 'C']),
                    chart_data.get('data', [10, 20, 30]),
                    chart_data.get('title', 'Data Analysis')
                )
            elif chart_type == 'pie':
                chart_path = self.graphics.create_pie_chart(
                    chart_data.get('data', [30, 30, 40]),
                    chart_data.get('labels', ['X', 'Y', 'Z']),
                    chart_data.get('title', 'Distribution')
                )
            elif chart_type == 'line':
                chart_path = self.graphics.create_line_chart(
                    chart_data.get('labels', [1, 2, 3, 4]),
                    chart_data.get('data', [5, 15, 10, 25]),
                    chart_data.get('title', 'Trend Analysis')
                )
            
            if chart_path and os.path.exists(chart_path):
                self.temp_files.append(chart_path)
                
                # Add picture to slide
                left = Inches(4)
                top = Inches(2)
                height = Inches(3)
                slide.shapes.add_picture(chart_path, left, top, height=height)
                
        except Exception as e:
            logger.error(f"Chart addition error: {e}")
    
    def _add_references_slide(self, references):
        """Add a references slide with citations"""
        if not references:
            return None
            
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        if title:
            title.text = "Source References"
        
        # References list
        content = slide.placeholders[1]
        if content:
            tf = content.text_frame
            tf.clear()
            
            for i, ref in enumerate(references[:10]):  # Limit to 10
                p = tf.add_paragraph()
                source = ref.get('source', 'Unknown')
                page = ref.get('page', '')
                page_str = f" (p. {page})" if page else ""
                p.text = f"{i+1}. {source}{page_str}"
                p.level = 0
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(0, 255, 0)
        
        return slide
    
    def export(self, slides_data, references=None, filename=None):
        """Export slides to PowerPoint with animations and graphics"""
        try:
            logger.info(f"Exporting {len(slides_data)} slides to PowerPoint")
            
            # Add slides
            for i, slide_data in enumerate(slides_data):
                if i == 0:
                    # First slide is title slide
                    self._add_title_slide(slide_data)
                else:
                    # Content slides
                    self._add_content_slide(slide_data, i)
            
            # Add references slide
            if references:
                self._add_references_slide(references)
            
            # Save to temp file
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                self.prs.save(tmp.name)
                tmp_path = tmp.name
            
            # Clean up temp files
            self._cleanup_temp_files()
            
            return tmp_path
            
        except Exception as e:
            logger.error(f"Export error: {e}")
            self._cleanup_temp_files()
            raise
    
    def add_vector_logo_to_all_slides(self):
        """Add Hawkins Lab logo to all slides"""
        for slide in self.prs.slides:
            self.graphics.add_vector_logo(slide, position=(9, 0.5), size=(0.8, 0.8))