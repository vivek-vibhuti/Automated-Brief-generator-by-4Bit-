"""
PowerPoint animation utilities for adding professional animations to slides
"""
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import xml.etree.ElementTree as ET

class AnimationManager:
    """Manages slide animations using PowerPoint's XML structure"""
    
    @staticmethod
    def add_fade_animation(shape, slide, timing_sec=0.5):
        """Add fade-in animation to a shape"""
        try:
            # Get or create animation extension list
            spTree = slide.shapes._spTree
            
            # Create animation element
            anim_elem = ET.SubElement(spTree, 'p:anim')
            anim_elem.set('to', '1')
            anim_elem.set('dur', f'{timing_sec * 1000}')
            anim_elem.set('autoRev', '0')
            
            # Create fade effect
            fade = ET.SubElement(anim_elem, 'p:cFade')
            fade.set('fade', 'in')
            
            return True
        except Exception as e:
            print(f"Animation error: {e}")
            return False
    
    @staticmethod
    def add_bullet_animations(slide, slide_data):
        """Add sequential fade animations to bullet points"""
        try:
            # Get content placeholder (usually index 1)
            if len(slide.shapes) > 1:
                content = slide.shapes[1]
                if content.has_text_frame:
                    for i, paragraph in enumerate(content.text_frame.paragraphs):
                        # Each bullet appears one by one
                        AnimationManager.add_fade_animation(
                            content, slide, timing_sec=0.3 + (i * 0.1)
                        )
            return True
        except Exception as e:
            print(f"Bullet animation error: {e}")
            return False
    
    @staticmethod
    def add_slide_transition(slide, transition_type='fade', duration=1.0):
        """Add transition to slide"""
        try:
            # Get slide XML
            slide_xml = slide.part.element
            
            # Add transition element
            transition = ET.SubElement(slide_xml, 'p:transition')
            transition.set('dur', f'{duration * 1000}')
            
            # Set transition type
            trans_type = ET.SubElement(transition, f'p:{transition_type}')
            
            return True
        except Exception as e:
            print(f"Transition error: {e}")
            return False

    @staticmethod
    def add_sequential_build(slide, elements_count):
        """Create sequential reveal animations"""
        for i in range(elements_count):
            ET.SubElement(slide.shapes._spTree, 'p:anim')
        return True