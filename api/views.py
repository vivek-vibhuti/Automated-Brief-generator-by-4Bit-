from django.shortcuts import render

# Create your views here.
from rest_framework import status
from rest_framework.decorators import api_view
from rest_framework.response import Response
from django.core.files.storage import default_storage
from django.core.files.base import ContentFile
import os
from django.http import FileResponse
from pptx import Presentation
from pptx.util import Inches, Pt
import json
import tempfile
import os

@api_view(['POST'])
def export_pptx(request):
    """Export slides to PowerPoint"""
    try:
        slides_data = request.data.get('slides', [])
        
        if not slides_data:
            return Response({'error': 'No slides data'}, status=400)
        
        # Create presentation
        prs = Presentation()
        
        for slide_data in slides_data:
            # Add slide with title and content layout
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data.get('title', 'Slide')
            
            # Set content (bullets)
            content = slide.placeholders[1].text_frame
            content.clear()
            
            bullets = slide_data.get('bullets', [])
            for i, bullet in enumerate(bullets):
                p = content.add_paragraph()
                p.text = bullet
                p.level = 0
                
                # Style
                p.font.size = Pt(18)
                if i == 0:
                    p.font.bold = True
            
            # Add sources as footer
            if slide_data.get('sources'):
                footer = slide.placeholders.get(2)  # Footer placeholder
                if footer:
                    footer.text = f"Sources: {', '.join(slide_data['sources'])}"
        
        # Save to temp file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            prs.save(tmp.name)
            tmp_path = tmp.name
        
        # Return file
        response = FileResponse(
            open(tmp_path, 'rb'),
            as_attachment=True,
            filename='hawkins_briefing.pptx'
        )
        
        # Clean up after sending
        import atexit
        atexit.register(lambda: os.unlink(tmp_path))
        
        return response
        
    except Exception as e:
        return Response({'error': str(e)}, status=500)

from .models import Document, DocumentChunk, Presentation, PresentationSource
from .serializers import DocumentSerializer, DocumentUploadSerializer, QuerySerializer, PresentationSerializer
from .utils.document_processor import DocumentProcessor
from .utils.chunking import SemanticChunker
from .utils.vector_store import VectorStore
from .utils.rag_retriever import RAGRetriever
from .utils.llm_clients import LLMClient
from .utils.slide_generator import RevealSlideGenerator
@api_view(['POST'])
def upload_document(request):
    """Upload and process a document"""
    try:
        print("\n" + "="*50)
        print("📁 UPLOAD DOCUMENT CALLED")
        print("="*50)
        print(f"Request FILES: {request.FILES}")
        print(f"Request data: {request.data}")
        
        serializer = DocumentUploadSerializer(data=request.data)
        if not serializer.is_valid():
            print(f"Serializer errors: {serializer.errors}")
            return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
        
        file = serializer.validated_data['file']
        title = serializer.validated_data.get('title', file.name)
        
        print(f"File: {file.name}, Size: {file.size}, Title: {title}")
        
        # Save document
        document = Document.objects.create(
            file=file,
            title=title
        )
        print(f"✅ Document saved with ID: {document.id}")
        
        # Process document
        try:
            processor = DocumentProcessor()
            chunks_data = processor.process_document(document)
            print(f"✅ Document processed: {len(chunks_data)} initial chunks")
            
            # Further chunking
            chunker = SemanticChunker()
            all_chunks = []
            for chunk in chunks_data:
                subchunks = chunker.chunk_text(
                    chunk['content'],
                    chunk['metadata']
                )
                all_chunks.extend(subchunks)
            
            print(f"✅ After semantic chunking: {len(all_chunks)} chunks")
            
            # Save chunks to DB
            db_chunks = []
            for i, chunk in enumerate(all_chunks):
                db_chunk = DocumentChunk.objects.create(
                    document=document,
                    content=chunk['content'],
                    page_number=chunk['metadata'].get('page_number'),
                    chunk_index=i,
                    metadata=chunk['metadata']
                )
                db_chunks.append(db_chunk)
            
            print(f"✅ Saved {len(db_chunks)} chunks to database")
            
            # Add to vector store
            try:
                vector_store = VectorStore()
                chunk_dicts = [{'content': c.content, 'metadata': c.metadata} for c in db_chunks]
                vector_store.add_chunks(chunk_dicts, document.id)
                print(f"✅ Added to vector store")
            except Exception as e:
                print(f"❌ Vector store error: {e}")
                import traceback
                traceback.print_exc()
                raise
            
            document.processed = True
            document.save()
            
            return Response({
                'message': 'Document uploaded and processed successfully',
                'document_id': document.id,
                'chunks_count': len(db_chunks)
            }, status=status.HTTP_201_CREATED)
            
        except Exception as e:
            print(f"❌ Processing error: {e}")
            import traceback
            traceback.print_exc()
            # Delete document if processing failed
            document.delete()
            return Response({
                'error': str(e),
                'type': e.__class__.__name__
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
        
    except Exception as e:
        print(f"❌ Upload error: {e}")
        import traceback
        traceback.print_exc()
        return Response({
            'error': str(e),
            'type': e.__class__.__name__
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
@api_view(['POST'])
def upload_document(request):
    """Upload and process a document"""
    serializer = DocumentUploadSerializer(data=request.data)
    if not serializer.is_valid():
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
    
    file = serializer.validated_data['file']
    title = serializer.validated_data.get('title', file.name)
    
    # Save document
    document = Document.objects.create(
        file=file,
        title=title
    )
    
    # Process document
    try:
        processor = DocumentProcessor()
        chunks_data = processor.process_document(document)
        
        # Further chunking
        chunker = SemanticChunker()
        all_chunks = []
        for chunk in chunks_data:
            subchunks = chunker.chunk_text(
                chunk['content'],
                chunk['metadata']
            )
            all_chunks.extend(subchunks)
        
        # Save chunks to DB
        db_chunks = []
        for i, chunk in enumerate(all_chunks):
            db_chunk = DocumentChunk.objects.create(
                document=document,
                content=chunk['content'],
                page_number=chunk['metadata'].get('page_number'),
                chunk_index=i,
                metadata=chunk['metadata']
            )
            db_chunks.append(db_chunk)
        
        # Add to vector store
        vector_store = VectorStore()
        chunk_dicts = [{'content': c.content, 'metadata': c.metadata} for c in db_chunks]
        vector_store.add_chunks(chunk_dicts, document.id)
        
        document.processed = True
        document.save()
        
        return Response({
            'message': 'Document uploaded and processed successfully',
            'document_id': document.id,
            'chunks_count': len(db_chunks)
        }, status=status.HTTP_201_CREATED)
        
    except Exception as e:
        return Response({
            'error': str(e)
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['GET'])
def list_documents(request):
    """List all uploaded documents"""
    documents = Document.objects.all().order_by('-uploaded_at')
    serializer = DocumentSerializer(documents, many=True)
    return Response(serializer.data)

@api_view(['POST'])
def generate_presentation(request):
    """Generate presentation from query"""
    serializer = QuerySerializer(data=request.data)
    if not serializer.is_valid():
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
    
    query = serializer.validated_data['query']
    document_ids = serializer.validated_data.get('document_ids')
    slide_count = serializer.validated_data.get('slide_count', 5)
    llm_provider = serializer.validated_data.get('llm_provider', 'gemini')
    
    try:
        # Retrieve relevant context
        retriever = RAGRetriever()
        context = retriever.retrieve_context(query, document_ids, top_k=10)
        references = retriever.get_references(query, document_ids, top_k=5)
        
        if not context:
            return Response({
                'error': 'No relevant documents found for your query'
            }, status=status.HTTP_404_NOT_FOUND)
        
        # Generate slide outline using LLM
        llm = LLMClient(provider=llm_provider)
        outline = llm.generate_slide_outline(query, context, slide_count)
        
        # Generate slides HTML
        slide_generator = RevealSlideGenerator()
        slides_html = slide_generator.generate(outline, references)
        
        # Save presentation
        presentation = Presentation.objects.create(
            query=query,
            slides_html=slides_html,
            slide_count=len(outline.get('slides', []))
        )
        
        # Save source references
        for ref in references:
            doc = Document.objects.filter(title=ref['source']).first()
            if doc:
                PresentationSource.objects.create(
                    presentation=presentation,
                    document=doc,
                    slide_index=0,  # You can improve this
                    relevance_score=ref.get('relevance', 0.0)
                )
        
        return Response({
            'presentation_id': presentation.id,
            'slides_html': slides_html,
            'slide_count': presentation.slide_count,
            'references': references
        })
        
    except Exception as e:
        return Response({
            'error': str(e)
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
@api_view(['POST'])
def export_pptx(request):
    """Export slides to PowerPoint"""
    try:
        # Log the request for debugging
        logger.info("PPTX export requested")
        
        # Get slides data from request
        data = request.data
        slides_data = data.get('slides', [])
        
        if not slides_data:
            logger.error("No slides data provided")
            return Response({'error': 'No slides data'}, status=400)
        
        logger.info(f"Exporting {len(slides_data)} slides to PowerPoint")
        
        # Create presentation
        prs = Presentation()
        
        for i, slide_data in enumerate(slides_data):
            logger.info(f"Processing slide {i+1}: {slide_data.get('title', 'Untitled')}")
            
            # Add slide with title and content layout
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            if title:
                title.text = slide_data.get('title', f'Slide {i+1}')
            
            # Set content (bullets)
            content = slide.placeholders[1]
            if content:
                text_frame = content.text_frame
                text_frame.clear()
                
                bullets = slide_data.get('bullets', [])
                for j, bullet in enumerate(bullets):
                    p = text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    
                    # Style
                    p.font.size = Pt(18)
                    if j == 0:
                        p.font.bold = True
            
            # Add sources as footer (if placeholder exists)
            if slide_data.get('sources') and len(slide.placeholders) > 2:
                footer = slide.placeholders[2]
                if footer:
                    footer.text = f"Sources: {', '.join(slide_data['sources'])}"
        
        # Save to temp file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            prs.save(tmp.name)
            tmp_path = tmp.name
            logger.info(f"PPTX saved to temporary file: {tmp_path}")
        
        # Return file
        response = FileResponse(
            open(tmp_path, 'rb'),
            as_attachment=True,
            filename=f'hawkins_briefing_{slides_data[0].get("title", "presentation")[:20]}.pptx'
        )
        
        # Clean up after sending
        import atexit
        atexit.register(lambda: os.unlink(tmp_path))
        
        return response
        
    except Exception as e:
        logger.error(f"PPTX export error: {str(e)}")
        import traceback
        traceback.print_exc()
        return Response({'error': str(e)}, status=500)
    from .utils.pptx_export import EnhancedPPTXExporter

@api_view(['POST'])
def export_pptx(request):
    """Export slides to PowerPoint with animations and vector graphics"""
    try:
        data = request.data
        slides_data = data.get('slides', [])
        references = data.get('references', [])
        
        if not slides_data:
            return Response({'error': 'No slides data'}, status=400)
        
        # Use enhanced exporter
        exporter = EnhancedPPTXExporter()
        pptx_path = exporter.export(slides_data, references)
        
        # Return file
        response = FileResponse(
            open(pptx_path, 'rb'),
            as_attachment=True,
            filename=f'hawkins_briefing_{len(slides_data)}slides.pptx'
        )
        
        # Clean up temp file after sending
        import atexit
        atexit.register(lambda: os.unlink(pptx_path))
        
        return response
        
    except Exception as e:
        logger.error(f"PPTX export error: {str(e)}")
        import traceback
        traceback.print_exc()
        return Response({'error': str(e)}, status=500)
@api_view(['GET'])
def get_presentation(request, presentation_id):
    """Get a generated presentation"""
    try:
        presentation = Presentation.objects.get(id=presentation_id)
        serializer = PresentationSerializer(presentation)
        return Response(serializer.data)
    except Presentation.DoesNotExist:
        return Response({
            'error': 'Presentation not found'
        }, status=status.HTTP_404_NOT_FOUND)

@api_view(['DELETE'])
def delete_document(request, document_id):
    """Delete a document and its chunks"""
    try:
        document = Document.objects.get(id=document_id)
        document.delete()
        # Note: Vector store cleanup would need additional handling
        return Response({
            'message': 'Document deleted successfully'
        })
    except Document.DoesNotExist:
        return Response({
            'error': 'Document not found'
        }, status=status.HTTP_404_NOT_FOUND)